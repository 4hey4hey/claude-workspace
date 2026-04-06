"""PPTX → Markdown 変換スクリプト

Usage:
    python scripts/pptx_to_md.py <input.pptx> <output.md>
"""

import sys
from pptx import Presentation
from pptx.util import Inches


def extract_table(table):
    """テーブルをMarkdown形式に変換"""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ''

    lines = []
    lines.append('| ' + ' | '.join(rows[0]) + ' |')
    lines.append('|' + '|'.join(['---' for _ in rows[0]]) + '|')
    for row in rows[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')
    return '\n'.join(lines)


def is_placeholder(shape):
    """プレースホルダーかどうかを安全に判定"""
    try:
        _ = shape.placeholder_format
        return True
    except ValueError:
        return False


def extract_slide(slide, index):
    """1スライドからMarkdownテキストを抽出"""
    slide_title = ''
    body_parts = []

    # まずタイトルプレースホルダーを探す
    for shape in slide.shapes:
        if is_placeholder(shape) and shape.placeholder_format.idx == 0:
            if shape.has_text_frame:
                slide_title = shape.text_frame.text.strip()
            break

    # 全テキストとテーブルを収集
    for shape in slide.shapes:
        if hasattr(shape, 'table'):
            body_parts.append(extract_table(shape.table))
        elif shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text != slide_title:
                    # インデントレベルに応じてリスト化
                    level = para.level if para.level else 0
                    if level > 0:
                        text = '  ' * (level - 1) + '- ' + text
                    body_parts.append(text)

    if not slide_title and body_parts:
        slide_title = body_parts[0]
        body_parts = body_parts[1:]

    lines = [f'## スライド {index}: {slide_title}', '']
    for part in body_parts:
        lines.append(part)
        lines.append('')
    lines.append('---')
    lines.append('')
    return '\n'.join(lines)


def pptx_to_md(input_path):
    """PPTXファイル全体をMarkdownに変換"""
    prs = Presentation(input_path)

    # メタ情報
    meta_lines = [
        f'# {input_path.split("/")[-1]}',
        '',
        f'**スライド数:** {len(prs.slides)}枚',
        '',
        '---',
        '',
    ]

    slide_lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines.append(extract_slide(slide, i))

    return '\n'.join(meta_lines + slide_lines)


def main():
    if len(sys.argv) < 3:
        print('Usage: python pptx_to_md.py <input.pptx> <output.md>')
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    md = pptx_to_md(input_path)

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(md)

    print(f'Converted: {input_path} → {output_path}')


if __name__ == '__main__':
    main()
